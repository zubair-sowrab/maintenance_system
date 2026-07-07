from pptx import Presentation
from pptx.util import Inches

# Initialize presentation
prs = Presentation()

def add_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        # Indent if the point starts with a dash
        p.level = 1 if point.startswith("-") else 0

# --- Slide 1: Title ---
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Real Estate Security Architecture"
title_slide.placeholders[1].text = "10-Layer Defense & PACT ERP Integration Specification"

# --- Slide 2: Layers 1 & 2 ---
add_slide("Layers 1 & 2: Server & Database Isolation", [
    "Layer 1: Port Control",
    "- EXPOSED: 443 (HTTPS) and 22 (SSH - restricted to you)",
    "- BLOCKED: 5432 (Postgres), 3306 (MySQL), 6379 (Redis), 8000 (Django)",
    "- Impact: External attacks cannot scan internal services",
    "Layer 2: PostgreSQL DB",
    "- Accepts connections from localhost only",
    "- Uses dedicated Django user with strict permissions",
    "- Impact: Unreachable from the outside internet"
])

# --- Slide 3: Layers 3 & 4 ---
add_slide("Layers 3 & 4: Application Core Defenses", [
    "Layer 3: Django Built-ins",
    "- SQL Injection: Prevented by mandatory ORM use",
    "- XSS: Blocked via automatic template escaping",
    "- CSRF & Clickjacking: Prevented natively",
    "Layer 4: Authentication",
    "- Email verification & strong password policies",
    "- Lockout enforced after 10 failed login attempts",
    "- Mandatory 2FA for Administrators"
])

# --- Slide 4: Django Image Placeholder ---
img_slide = prs.slides.add_slide(prs.slide_layouts[5])
img_slide.shapes.title.text = "Layer 3 Detail: Django Protections"
txBox = img_slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
txBox.text_frame.text = "[ INSERT YOUR DJANGO SECURITY DIAGRAM IMAGE HERE ]"

# --- Slide 5: Layers 5 & 6 ---
add_slide("Layers 5 & 6: Permissions & Transit", [
    "Layer 5: Authorization",
    "- Strict ID matching for Tenants and Owners",
    "- Role-based permissions separate views",
    "- Impact: Tenants cannot modify URLs to view other leases",
    "Layer 6: HTTPS Encryption",
    "- 100% traffic encrypted between User and Django",
    "- 100% traffic encrypted between Python Sync and PACT"
])

# --- Slide 6: Layers 7 & 8 ---
add_slide("Layers 7 & 8: Traffic & Edge Control", [
    "Layer 7: Rate Limiting",
    "- Slows down rapid, unnatural request spikes",
    "- Blacklists suspicious user IPs dynamically",
    "Layer 8: Firewalls",
    "- Digital Ocean Cloud Firewall (Outer boundary)",
    "- Linux Firewall - UFW/nftables (Inner boundary)",
    "- Impact: Drops unauthorized connections immediately"
])

# --- Slide 7: Layers 9 & 10 ---
add_slide("Layers 9 & 10: Perimeter & Assets", [
    "Layer 9: DDoS Protection",
    "- Reverse Proxy / CDN (Cloudflare) in front of server",
    "- Impact: Absorbs massive traffic floods",
    "Layer 10: Secure File Uploads",
    "- Uploaded documents stored outside the web root (e.g., S3)",
    "- Impact: Prevents execution of malicious virus uploads"
])

# --- Slide 8: PACT Rules ---
add_slide("Zero-Trust: PACT ERP Rules", [
    "Total Isolation: Users do not know PACT exists",
    "No Keys in Django: Django app does not hold the API key",
    "- Only Python Sync holds the credentials",
    "IP Restriction: PACT firewall allows ONLY the DO Server IP",
    "Read-Only Access: Only GET Methods allowed",
    "- NO Delete, Update, Create, or Insert commands"
])

# Save the file
prs.save("Real_Estate_Security_Deck.pptx")
print("Successfully generated: Real_Estate_Security_Deck.pptx")